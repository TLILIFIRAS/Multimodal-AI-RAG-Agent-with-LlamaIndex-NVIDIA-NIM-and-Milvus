import os
import fitz
from pptx import Presentation
import subprocess
from llama_index.core import Document
from utils import (
    describe_image, is_graph, process_graph, extract_text_around_item, 
    process_text_blocks, save_uploaded_file
)


def get_pdf_documents(pdf_file):
    """
    Process a PDF file and extract text, tables, and images.

    Args:
        pdf_file (UploadedFile): The PDF file object to be processed.

    Returns:
        all_pdf_documents (list): A list of `Document` objects containing extracted text, tables, and images.
    """
    
    all_pdf_documents = []  # List to store all extracted documents (text, tables, images).
    ongoing_tables = {}  # Dictionary to track ongoing tables during processing.

    try:
        # Open the PDF file stream.
        f = fitz.open(stream=pdf_file.read(), filetype="pdf")
    except Exception as e:
        # Handle errors during PDF processing and return an empty list.
        print(f"Error opening or processing the PDF file: {e}")
        return []

    # Iterate over each page in the PDF.
    for i in range(len(f)):
        page = f[i]
        
        # Extract text blocks from the page, filtering out the ones outside the main content area.
        text_blocks = [block for block in page.get_text("blocks", sort=True) 
                       if block[-1] == 0 and not (block[1] < page.rect.height * 0.1 or block[3] > page.rect.height * 0.9)]
        
        # Group the extracted text blocks into chunks based on character count.
        grouped_text_blocks = process_text_blocks(text_blocks)
        
        # Parse tables from the page and track ongoing tables.
        table_docs, table_bboxes, ongoing_tables = parse_all_tables(pdf_file.name, page, i, text_blocks, ongoing_tables)
        
        # Add extracted table documents to the result list.
        all_pdf_documents.extend(table_docs)

        # Parse images from the page.
        image_docs = parse_all_images(pdf_file.name, page, i, text_blocks)
        
        # Add extracted image documents to the result list.
        all_pdf_documents.extend(image_docs)

        # Process each grouped text block and create a document for it.
        for text_block_ctr, (heading_block, content) in enumerate(grouped_text_blocks, 1):
            heading_bbox = fitz.Rect(heading_block[:4])
            
            # Skip any heading block that intersects with a table.
            if not any(heading_bbox.intersects(table_bbox) for table_bbox in table_bboxes):
                # Create a bounding box dictionary for the text block.
                bbox = {"x1": heading_block[0], "y1": heading_block[1], "x2": heading_block[2], "x3": heading_block[3]}
                
                # Create a document for the text block and add it to the list of extracted documents.
                text_doc = Document(
                    text=f"{heading_block[4]}\n{content}",
                    metadata={  # Include metadata such as bounding box, page number, and source information.
                        **bbox,
                        "type": "text",
                        "page_num": i,
                        "source": f"{pdf_file.name[:-4]}-page{i}-block{text_block_ctr}"
                    },
                    id_=f"{pdf_file.name[:-4]}-page{i}-block{text_block_ctr}"  # Unique ID for the document.
                )
                all_pdf_documents.append(text_doc)

    # Close the PDF file after processing.
    f.close()

    # Return the list of all extracted documents.
    return all_pdf_documents

def parse_all_tables(filename, page, pagenum, text_blocks, ongoing_tables):
    """
    Extract tables from a PDF page.

    Args:
        filename (str): The name of the PDF file.
        page (fitz.Page): The current page of the PDF.
        pagenum (int): The page number being processed.
        text_blocks (list): List of text blocks extracted from the page.
        ongoing_tables (dict): A dictionary that tracks ongoing tables (not used in this function, but could be useful for state management).

    Returns:
        table_docs (list): List of documents containing the extracted tables.
        table_bboxes (list): List of bounding boxes for the extracted tables.
        ongoing_tables (dict): The updated ongoing tables dictionary (currently not modified in this function).
    """
    
    table_docs = []  # List to store table documents.
    table_bboxes = []  # List to store the bounding boxes of the tables.

    try:
        # Find tables on the page using strict horizontal and vertical strategies.
        tables = page.find_tables(horizontal_strategy="lines_strict", vertical_strategy="lines_strict")
        
        # Iterate through each table found on the page.
        for tab in tables:
            if not tab.header.external:  # Check if the table has no external header.
                # Convert the table to a Pandas DataFrame.
                pandas_df = tab.to_pandas()
                
                # Create a directory to store table references (Excel files and images).
                tablerefdir = os.path.join(os.getcwd(), "vectorstore/table_references")
                os.makedirs(tablerefdir, exist_ok=True)

                # Save the DataFrame as an Excel file.
                df_xlsx_path = os.path.join(tablerefdir, f"table{len(table_docs)+1}-page{pagenum}.xlsx")
                pandas_df.to_excel(df_xlsx_path)

                # Create a bounding box for the table.
                bbox = fitz.Rect(tab.bbox)
                table_bboxes.append(bbox)

                # Extract surrounding text above and below the table (for captioning purposes).
                before_text, after_text = extract_text_around_item(text_blocks, bbox, page.rect.height)

                # Capture an image of the table using the bounding box.
                table_img = page.get_pixmap(clip=bbox)
                table_img_path = os.path.join(tablerefdir, f"table{len(table_docs)+1}-page{pagenum}.jpg")
                table_img.save(table_img_path)

                # Process the table image (e.g., generate a description).
                description = process_graph(table_img.tobytes())

                # Combine the surrounding text with the description to create a caption.
                caption = before_text.replace("\n", " ") + description + after_text.replace("\n", " ")
                
                # If no surrounding text exists, use the table's header names as the caption.
                if before_text == "" and after_text == "":
                    caption = " ".join(tab.header.names)
                
                # Define metadata for the table document.
                table_metadata = {
                    "source": f"{filename[:-4]}-page{pagenum}-table{len(table_docs)+1}",
                    "dataframe": df_xlsx_path,
                    "image": table_img_path,
                    "caption": caption,
                    "type": "table",
                    "page_num": pagenum
                }

                # Prepare the text to describe the table's contents (including column names).
                all_cols = ", ".join(list(pandas_df.columns.values))
                doc = Document(text=f"This is a table with the caption: {caption}\nThe columns are {all_cols}", metadata=table_metadata)

                # Append the table document to the result list.
                table_docs.append(doc)

    except Exception as e:
        print(f"Error during table extraction: {e}")

    # Return the extracted table documents, bounding boxes, and the ongoing tables dictionary.
    return table_docs, table_bboxes, ongoing_tables

def parse_all_images(filename, page, pagenum, text_blocks):
    """
    Extract images from a PDF page.

    Args:
        filename (str): The name of the PDF file.
        page (fitz.Page): The current page of the PDF.
        pagenum (int): The page number being processed.
        text_blocks (list): List of text blocks extracted from the page.

    Returns:
        image_docs (list): List of documents containing the extracted images and their metadata.
    """
    
    image_docs = []  # List to store image documents.
    
    # Get information about all images on the current page, including their xrefs.
    image_info_list = page.get_image_info(xrefs=True)
    
    # Get the dimensions of the page.
    page_rect = page.rect

    # Iterate over all images on the page.
    for image_info in image_info_list:
        xref = image_info['xref']
        
        # Skip invalid image references (xref 0 means no image found).
        if xref == 0:
            continue

        # Define the bounding box of the image from its information.
        img_bbox = fitz.Rect(image_info['bbox'])

        # Skip small images (less than 5% of the page's width or height).
        if img_bbox.width < page_rect.width / 20 or img_bbox.height < page_rect.height / 20:
            continue

        # Extract the image data.
        extracted_image = page.parent.extract_image(xref)
        image_data = extracted_image["image"]

        # Create a directory to store the image references (PNG format).
        imgrefpath = os.path.join(os.getcwd(), "vectorstore/image_references")
        os.makedirs(imgrefpath, exist_ok=True)

        # Save the extracted image as a PNG file.
        image_path = os.path.join(imgrefpath, f"image{xref}-page{pagenum}.png")
        with open(image_path, "wb") as img_file:
            img_file.write(image_data)

        # Extract the surrounding text before and after the image (for captioning purposes).
        before_text, after_text = extract_text_around_item(text_blocks, img_bbox, page.rect.height)
        
        # Skip if no surrounding text is found.
        if before_text == "" and after_text == "":
            continue

        # Process the image if it's a graph or chart (add description).
        image_description = " "
        if is_graph(image_data):
            image_description = process_graph(image_data)

        # Combine the surrounding text with the description to form a caption.
        caption = before_text.replace("\n", " ") + image_description + after_text.replace("\n", " ")

        # Define the metadata for the image document.
        image_metadata = {
            "source": f"{filename[:-4]}-page{pagenum}-image{xref}",
            "image": image_path,
            "caption": caption,
            "type": "image",
            "page_num": pagenum
        }

        # Create the image document and append it to the list.
        image_docs.append(Document(text="This is an image with the caption: " + caption, metadata=image_metadata))

    # Return the list of image documents.
    return image_docs

def process_ppt_file(ppt_path):
    """Process a PowerPoint file."""
    
    # Step 1: Convert PowerPoint to PDF
    pdf_path = convert_ppt_to_pdf(ppt_path)
    
    # Step 2: Convert PDF to images (slides converted into image format)
    images_data = convert_pdf_to_images(pdf_path)
    
    # Step 3: Extract text and speaker notes from the PowerPoint file
    slide_texts = extract_text_and_notes_from_ppt(ppt_path)
    
    # Step 4: Initialize an empty list to store the processed data (documents)
    processed_data = []

    # Step 5: Process each slide's data (image, text, and notes)
    for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
        
        # Step 5.1: If there are speaker notes, prepend a note text
        if notes:
            notes = "\n\nThe speaker notes for this slide are: " + notes
        
        # Step 5.2: Read the image file in binary mode
        with open(image_path, 'rb') as image_file:
            image_content = image_file.read()
        
        # Step 5.3: If the image is a graph, process it to generate a description
        image_description = " "
        if is_graph(image_content):
            image_description = process_graph(image_content)
        
        # Step 5.4: Prepare metadata for the image document
        image_metadata = {
            "source": f"{os.path.basename(ppt_path)}",  # File name (without path)
            "image": image_path,  # Path to the image file
            "caption": slide_text + image_description + notes,  # Caption combining text, description, and notes
            "type": "image",  # Type of content
            "page_num": page_num  # Slide number (page number in PDF)
        }

        # Step 5.5: Create a Document object containing the processed content
        processed_data.append(Document(text="This is a slide with the text: " + slide_text + image_description, metadata=image_metadata))

    # Step 6: Return the list of processed documents
    return processed_data

def convert_ppt_to_pdf(ppt_path):
    """Convert a PowerPoint file to PDF using LibreOffice."""
    
    # Step 1: Extract the base name and remove file extension
    base_name = os.path.basename(ppt_path)  # Get the base file name from the path (e.g., 'presentation.pptx')
    ppt_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')  # Remove spaces and file extension
    
    # Step 2: Define the directory path where the PDF will be saved
    new_dir_path = os.path.abspath("vectorstore/ppt_references")  # Directory where converted PDF will be stored
    os.makedirs(new_dir_path, exist_ok=True)  # Create the directory if it does not exist
    
    # Step 3: Define the output PDF file path
    pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")  # Final path for the converted PDF
    
    # Step 4: Run LibreOffice to convert the PPT to PDF
    command = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', new_dir_path, ppt_path]
    subprocess.run(command, check=True)  # Run the command and check for errors
    
    # Step 5: Return the path of the converted PDF file
    return pdf_path

def convert_pdf_to_images(pdf_path):
    """Convert a PDF file to a series of images using PyMuPDF."""
    
    # Step 1: Open the PDF file using PyMuPDF
    doc = fitz.open(pdf_path)
    
    # Step 2: Get the base name of the PDF and generate a clean name without spaces
    base_name = os.path.basename(pdf_path)  # Get the filename with extension (e.g., 'document.pdf')
    pdf_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')  # Remove extension and spaces
    
    # Step 3: Define the output directory where images will be saved
    new_dir_path = os.path.join(os.getcwd(), "vectorstore/ppt_references")  # Path for saving images
    os.makedirs(new_dir_path, exist_ok=True)  # Create the directory if it does not exist
    
    # Step 4: Initialize an empty list to store image paths
    image_paths = []

    # Step 5: Loop through each page of the PDF and convert to image
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)  # Load the page
        pix = page.get_pixmap()  # Convert the page to a pixmap (image)
        
        # Step 6: Save the image to the output directory
        output_image_path = os.path.join(new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png")  # File path
        pix.save(output_image_path)  # Save the image as PNG
        
        # Step 7: Append the image path and page number as a tuple to the list
        image_paths.append((output_image_path, page_num))

    # Step 8: Close the PDF document
    doc.close()

    # Step 9: Return the list of image paths and corresponding page numbers
    return image_paths

def extract_text_and_notes_from_ppt(ppt_path):
    """Extract text and notes from a PowerPoint file."""
    
    # Step 1: Open the PowerPoint file using the python-pptx library
    prs = Presentation(ppt_path)  # Load the PowerPoint presentation

    # Step 2: Initialize an empty list to store text and notes from each slide
    text_and_notes = []

    # Step 3: Iterate over each slide in the presentation
    for slide in prs.slides:
        
        # Step 4: Extract the text from each shape in the slide (e.g., textboxes, titles, etc.)
        slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        
        # Step 5: Extract the notes for the slide (if available)
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
        except:
            notes = ''
        
        # Step 6: Append the extracted text and notes as a tuple to the list
        text_and_notes.append((slide_text, notes))

    # Step 7: Return the list of text and notes from each slide
    return text_and_notes

def load_multimodal_data(files):
    """Load and process multiple file types."""
    documents = []  # List to store processed documents
    
    for file in files:  # Iterate through each uploaded file
        file_extension = os.path.splitext(file.name.lower())[1]  # Get the file extension in lowercase
        
        if file_extension in ('.png', '.jpg', '.jpeg'):  # If the file is an image
            image_content = file.read()  # Read the image file
            image_text = describe_image(image_content)  # Use a function (e.g., describe_image) to extract a description from the image
            doc = Document(text=image_text, metadata={"source": file.name, "type": "image"})  # Create a Document object
            documents.append(doc)  # Add the Document to the list
            
        elif file_extension == '.pdf':  # If the file is a PDF
            try:
                pdf_documents = get_pdf_documents(file)  # Process the PDF file and extract text/content
                documents.extend(pdf_documents)  # Add the processed PDF documents to the list
            except Exception as e:  # Handle errors in PDF processing
                print(f"Error processing PDF {file.name}: {e}")
                
        elif file_extension in ('.ppt', '.pptx'):  # If the file is a PowerPoint
            try:
                ppt_documents = process_ppt_file(save_uploaded_file(file))  # Process the PowerPoint file
                documents.extend(ppt_documents)  # Add the processed PowerPoint documents to the list
            except Exception as e:  # Handle errors in PowerPoint processing
                print(f"Error processing PPT {file.name}: {e}")
                
        else:  # If the file is a text file (or another type)
            text = file.read().decode("utf-8")  # Read the text content
            doc = Document(text=text, metadata={"source": file.name, "type": "text"})  # Create a Document object
            documents.append(doc)  # Add the Document to the list
    
    return documents  # Return the list of processed documents


def load_data_from_directory(directory):
    """Load and process multiple file types from a directory."""
    documents = []  # List to store processed documents
    
    # Iterate over all files in the specified directory
    for filename in os.listdir(directory):
        filepath = os.path.join(directory, filename)  # Get full path of the file
        file_extension = os.path.splitext(filename.lower())[1]  # Get file extension (e.g., .png, .pdf)
        
        print(filename)  # Print the filename being processed
        
        # If the file is an image (PNG, JPG, JPEG)
        if file_extension in ('.png', '.jpg', '.jpeg'):
            with open(filepath, "rb") as image_file:  # Open the image file in binary read mode
                image_content = image_file.read()  # Read the image content
            image_text = describe_image(image_content)  # Extract a textual description of the image
            doc = Document(text=image_text, metadata={"source": filename, "type": "image"})  # Create a Document object for the image
            print(doc)  # Print the Document object (for debugging purposes)
            documents.append(doc)  # Add the image document to the list of documents
        
        # If the file is a PDF
        elif file_extension == '.pdf':
            with open(filepath, "rb") as pdf_file:  # Open the PDF file in binary read mode
                try:
                    pdf_documents = get_pdf_documents(pdf_file)  # Process the PDF and extract documents
                    documents.extend(pdf_documents)  # Add the extracted documents to the list
                except Exception as e:
                    print(f"Error processing PDF {filename}: {e}")  # Handle any errors during PDF processing
        
        # If the file is a PowerPoint (PPT or PPTX)
        elif file_extension in ('.ppt', '.pptx'):
            try:
                ppt_documents = process_ppt_file(filepath)  # Process the PowerPoint file and extract documents
                documents.extend(ppt_documents)  # Add the extracted documents to the list
                print(ppt_documents)  # Print the extracted PowerPoint documents (for debugging purposes)
            except Exception as e:
                print(f"Error processing PPT {filename}: {e}")  # Handle any errors during PowerPoint processing
        
        # If the file is a text file
        else:
            with open(filepath, "r", encoding="utf-8") as text_file:  # Open the text file with UTF-8 encoding
                text = text_file.read()  # Read the content of the text file
            doc = Document(text=text, metadata={"source": filename, "type": "text"})  # Create a Document object for the text file
            documents.append(doc)  # Add the text document to the list of documents
    
    return documents  # Return the list of processed documents
